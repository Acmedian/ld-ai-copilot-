import os
import json
import base64
import streamlit as st
from openai import OpenAI
from audiorecorder import audiorecorder
from pptx import Presentation

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

st.set_page_config(page_title="AI L&D Copilot", layout="wide")

st.title("🎤 AI Learning Design Copilot")
st.caption("Voice conversation → Training Design → Slide Generation")

st.divider()

# -------------------------------------------------
# SESSION STATE
# -------------------------------------------------

if "messages" not in st.session_state:
    st.session_state.messages = []

if "design" not in st.session_state:
    st.session_state.design = {}

if "slides" not in st.session_state:
    st.session_state.slides = []

if "spoken_index" not in st.session_state:
    st.session_state.spoken_index = -1


# -------------------------------------------------
# PAGE LAYOUT
# -------------------------------------------------

col1, col2, col3 = st.columns([1.2,1,1.3])


# =================================================
# COLUMN 1 — VOICE CONVERSATION
# =================================================

with col1:

    st.subheader("🎙 SME Conversation")

    audio = audiorecorder("Speak", "Recording...")

    if len(audio) > 0:

        audio.export("input.wav", format="wav")

        with open("input.wav", "rb") as f:

            transcript = client.audio.transcriptions.create(
                model="gpt-4o-mini-transcribe",
                file=f
            )

        user_text = transcript.text

        st.session_state.messages.append({
            "role":"user",
            "content":user_text
        })

        with st.chat_message("user"):
            st.write(user_text)

        system_prompt = """
You are an expert learning designer.

Interview the SME to gather:
• training title
• target audience
• learning objectives
• modules
• assessment

Ask ONE short question at a time.
Keep responses under 2 sentences.
"""

        messages = [{"role":"system","content":system_prompt}]
        messages += st.session_state.messages

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages
        )

        ai_text = response.choices[0].message.content

        st.session_state.messages.append({
            "role":"assistant",
            "content":ai_text
        })

    assistant_count = 0

    for msg in st.session_state.messages:

        if msg["role"] == "assistant":

            with st.chat_message("assistant"):
                st.write(msg["content"])

            if assistant_count > st.session_state.spoken_index:

                st.session_state.spoken_index = assistant_count

                speech = client.audio.speech.create(
                    model="gpt-4o-mini-tts",
                    voice="alloy",
                    input=msg["content"][:200]
                )

                with open("response.wav","wb") as f:
                    f.write(speech.content)

                with open("response.wav","rb") as f:
                    audio_bytes = f.read()

                b64 = base64.b64encode(audio_bytes).decode()

                st.markdown(
                    f"""
                    <audio autoplay>
                    <source src="data:audio/wav;base64,{b64}" type="audio/wav">
                    </audio>
                    """,
                    unsafe_allow_html=True
                )

            assistant_count += 1


# =================================================
# COLUMN 2 — TRAINING DESIGN
# =================================================

with col2:

    st.subheader("📋 Training Blueprint")

    if len(st.session_state.messages) > 0:

        design_response = client.chat.completions.create(
            model="gpt-4o-mini",
            response_format={"type":"json_object"},
            messages=[
                {
                    "role":"system",
                    "content":"""
Extract training design.

Return JSON with:

title
audience
objectives
modules
assessment
"""
                },
                {
                    "role":"user",
                    "content":str(st.session_state.messages)
                }
            ]
        )

        design_data = json.loads(design_response.choices[0].message.content)

        st.session_state.design = design_data

    design = st.session_state.design

    st.markdown("### Training Title")
    st.write(design.get("title",""))

    st.markdown("### Audience")
    st.write(design.get("audience",""))

    st.markdown("### Objectives")

    for obj in design.get("objectives",[]):
        st.write("•",obj)

    st.markdown("### Modules")

    for mod in design.get("modules",[]):
        st.write("•",mod)

    st.markdown("### Assessment")
    st.write(design.get("assessment",""))


# =================================================
# COLUMN 3 — VISUAL SLIDES
# =================================================

with col3:

    st.subheader("🎨 Slide Preview")

    if st.session_state.design:

        slide_response = client.chat.completions.create(
            model="gpt-4o-mini",
            response_format={"type":"json_object"},
            messages=[
                {
                    "role":"system",
                    "content":"""
Create presentation slides.

Return JSON:

slides:[
{
title:"",
bullets:["","",""]
}
]
"""
                },
                {
                    "role":"user",
                    "content":json.dumps(st.session_state.design)
                }
            ]
        )

        slide_data = json.loads(slide_response.choices[0].message.content)

        slides = slide_data.get("slides",[])

        st.session_state.slides = slides

        slide_number = 1

        for slide in slides:

            bullets_html = "".join(
                [f"<li>{b}</li>" for b in slide["bullets"]]
            )

            slide_html = f"""
            <div style="
                background:white;
                border-radius:14px;
                overflow:hidden;
                margin-bottom:30px;
                box-shadow:0 8px 20px rgba(0,0,0,0.15);
                border:1px solid #e5e5e5;
            ">

            <div style="
                background:#C00000;
                color:white;
                padding:16px;
                font-size:20px;
                font-weight:600;
            ">
            {slide["title"]}
            </div>

            <div style="
                padding:22px;
                font-size:16px;
                line-height:1.7;
            ">

            <div style="
                font-size:12px;
                color:#888;
                margin-bottom:12px;
            ">
            Slide {slide_number}
            </div>

            <ul>
            {bullets_html}
            </ul>

            </div>

            </div>
            """

            st.markdown(slide_html, unsafe_allow_html=True)

            slide_number += 1


    # ------------------------------------------------
    # POWERPOINT DOWNLOAD
    # ------------------------------------------------

    if st.session_state.slides:

        prs = Presentation()

        for slide in st.session_state.slides:

            layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)

            s.shapes.title.text = slide["title"]

            content = "\n".join(slide["bullets"])

            s.placeholders[1].text = content

        prs.save("training_deck.pptx")

        with open("training_deck.pptx","rb") as f:

            st.download_button(
                "⬇ Download PowerPoint",
                f,
                file_name="training_deck.pptx"
            )
